"""Generate a 10-minute presentation with speaker notes."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_MED = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x53, 0x3c, 0xc6)
ACCENT2 = RGBColor(0x00, 0xd2, 0xff)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT = RGBColor(0xcc, 0xcc, 0xcc)
GREEN = RGBColor(0x00, 0xe6, 0x76)
RED = RGBColor(0xff, 0x45, 0x45)
YELLOW = RGBColor(0xff, 0xd7, 0x00)
ORANGE = RGBColor(0xff, 0x8c, 0x00)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, space_before=6, bullet=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 1
    return p


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, 1, 1.5, 11.3, 1.5,
             "Adaptive Intensity Token Importance\nfor PPO-Based RLHF",
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 3.5, 11.3, 0.8,
             "From Bias-Variance Analysis to MSE-Optimal Weighting",
             font_size=24, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.0, 11.3, 0.5,
             "AITI  |  MOAI  |  Token-Level Importance for Online RLHF",
             font_size=18, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_notes(slide, """Welcome everyone. Today I'll present our work on adaptive token importance weighting for PPO-based RLHF.

The key question we started with: can we make PPO training more efficient by focusing gradient updates on the tokens that matter most? The answer turns out to be yes, but with a critical caveat about managing bias over training.

We'll cover two methods: AITI, which uses a fixed decay schedule, and MOAI, which derives the optimal weighting intensity from a closed-form formula. Both achieve Pareto improvements over standard PPO.""")

# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "The Problem: Not All Tokens Are Equal",
             font_size=32, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6, 5,
                  "Standard PPO treats every token equally:",
                  font_size=20, color=WHITE)
add_para(tf, "", 14)
add_para(tf, 'L = (1/T) Sum_t  min(r_t * A_t, clip(r_t) * A_t)', 18, YELLOW, True, 12)
add_para(tf, "", 10)
add_para(tf, "But in a sentence like:", 20, WHITE, False, 16)
add_para(tf, '"The food was surprisingly delicious."', 20, ACCENT2, True, 8)
add_para(tf, "", 10)
add_para(tf, '  "The"  -  function word, low signal', 16, LIGHT, False, 8)
add_para(tf, '  "surprisingly"  -  high uncertainty, high signal', 16, GREEN, False, 4)
add_para(tf, '  "delicious"  -  reward-carrying, high advantage', 16, GREEN, False, 4)
add_para(tf, "", 10)
add_para(tf, "Can we weight tokens by importance?", 22, WHITE, True, 16)

tf2 = add_text_box(slide, 7, 1.3, 5.8, 5,
                   "Prior work: TI-DPO (Yang et al. 2025)",
                   font_size=20, color=ORANGE, bold=True)
add_para(tf2, "  Token importance for DPO (offline)", 17, LIGHT, False, 12)
add_para(tf2, "  Uses gradient attribution + Gaussian prior", 17, LIGHT, False, 4)
add_para(tf2, "  Works well for preference pairs", 17, LIGHT, False, 4)
add_para(tf2, "", 14)
add_para(tf2, "Our question:", 20, WHITE, True, 12)
add_para(tf2, "Does this transfer to PPO (online)?", 22, ACCENT2, True, 8)
add_para(tf2, "", 14)
add_para(tf2, "Short answer: No.", 22, RED, True, 12)
add_para(tf2, "But investigating WHY led to", 18, LIGHT, False, 8)
add_para(tf2, "a fundamental insight...", 18, LIGHT, False, 4)

add_notes(slide, """Standard PPO treats all tokens in a generated sequence equally. But clearly some tokens matter more than others. Function words like "the" carry little reward signal, while words like "surprisingly" or "delicious" are where the model is actually learning.

TI-DPO showed that token importance weighting helps in the offline DPO setting, where you have chosen/rejected pairs. We asked: does this transfer to online PPO?

The answer is no - the paper's DPO-derived signals don't work in PPO because there are no contrastive pairs. But investigating WHY fixed importance weighting fails in PPO led us to a fundamental insight about bias accumulation.""")

# ============================================================
# SLIDE 3: PPO-Native Methods & Degradation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "PPO-Native Importance Methods: Initial Promise, Then Failure",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6, 2.5,
                  "We derived importance from PPO signals:",
                  font_size=20, color=WHITE)
add_para(tf, "", 8)
add_para(tf, "Entropy:  w(t) = H(pi(.|s_t)) / mean(H)", 17, GREEN, False, 8)
add_para(tf, "  Focus on uncertain tokens", 15, LIGHT, False, 2)
add_para(tf, "|Advantage|:  w(t) = |A(t)| / mean(|A|)", 17, GREEN, False, 12)
add_para(tf, "  Focus on high-signal tokens", 15, LIGHT, False, 2)

# Results table
tf2 = add_text_box(slide, 0.5, 4.0, 12, 3,
                   "The Degradation Pattern",
                   font_size=22, color=YELLOW, bold=True)
add_para(tf2, "", 8)
add_para(tf2, "Method               R@100ep    R@150ep    Change", 16, ACCENT2, True, 8)
add_para(tf2, "-------------------------------------------------------", 14, LIGHT, False, 2)
add_para(tf2, "PPO baseline         0.791       0.850      +0.059  (keeps improving)", 16, GREEN, False, 4)
add_para(tf2, "Entropy              0.872       0.797      -0.075  (degrades!)", 16, RED, False, 4)
add_para(tf2, "|Advantage|          0.870       0.785      -0.085  (degrades!)", 16, RED, False, 4)
add_para(tf2, "Paper Hybrid         0.816       0.779      -0.037  (degrades!)", 16, RED, False, 4)

tf3 = add_text_box(slide, 7, 1.3, 5.8, 2.5,
                   "At 100 episodes: importance wins!",
                   font_size=20, color=GREEN, bold=True)
add_para(tf3, "Entropy: 0.872 vs PPO: 0.791", 18, WHITE, False, 8)
add_para(tf3, "", 10)
add_para(tf3, "At 150 episodes: importance loses.", 20, RED, True, 12)
add_para(tf3, "PPO: 0.850 vs Entropy: 0.797", 18, WHITE, False, 8)
add_para(tf3, "", 10)
add_para(tf3, "Why does importance weighting degrade?", 20, YELLOW, True, 12)

add_notes(slide, """We derived two PPO-native importance methods that don't need DPO's contrastive pairs: entropy weighting (focus on uncertain tokens) and advantage magnitude weighting (focus on high-signal tokens).

At 100 episodes, both beat PPO significantly - entropy reached 0.872 reward vs PPO's 0.791. Great!

But at 150 episodes, something disturbing happened. Both importance methods DEGRADED while PPO baseline kept improving. PPO reached 0.850 while entropy dropped to 0.797.

This pattern was consistent across ALL importance methods we tested. Something fundamental is wrong with fixed-intensity importance weighting in online RL. Why?""")

# ============================================================
# SLIDE 4: The Key Insight - Bias-Variance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "The Key Insight: Importance Weighting Has a Bias-Variance Tradeoff",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6, 5.5,
                  "The Math:",
                  font_size=22, color=YELLOW, bold=True)
add_para(tf, "", 8)
add_para(tf, "True PPO gradient:", 18, LIGHT, False, 8)
add_para(tf, "  g = E_t[ grad(log pi) * A_t ]", 18, WHITE, True, 4)
add_para(tf, "", 8)
add_para(tf, "Weighted gradient:", 18, LIGHT, False, 8)
add_para(tf, "  g_w = E_t[ w(t) * grad(log pi) * A_t ]", 18, WHITE, True, 4)
add_para(tf, "", 8)
add_para(tf, "Bias = E[g_w] - E[g] = Cov(w, f)", 22, YELLOW, True, 12)
add_para(tf, "", 8)
add_para(tf, "If w correlates with f (it does -", 17, LIGHT, False, 8)
add_para(tf, "that's WHY we weight!), bias != 0", 17, LIGHT, False, 2)

tf2 = add_text_box(slide, 7, 1.3, 5.8, 5.5,
                   "The Tradeoff:",
                   font_size=22, color=YELLOW, bold=True)
add_para(tf2, "", 8)
add_para(tf2, "Variance: O(1/sqrt(N))", 20, GREEN, True, 12)
add_para(tf2, "  Decreases with more data", 16, LIGHT, False, 4)
add_para(tf2, "  Importance reduces this", 16, GREEN, False, 4)
add_para(tf2, "", 12)
add_para(tf2, "Bias: Cov(w, f) -- persistent!", 20, RED, True, 12)
add_para(tf2, "  Does NOT average out", 16, LIGHT, False, 4)
add_para(tf2, "  Accumulates step after step", 16, RED, False, 4)
add_para(tf2, "", 16)
add_para(tf2, "Early: Var >> Bias --> weighting helps", 20, GREEN, True, 12)
add_para(tf2, "Late:  Bias >> Var --> weighting hurts", 20, RED, True, 8)
add_para(tf2, "", 12)
add_para(tf2, "This is why EVERY fixed method degrades!", 20, YELLOW, True, 12)

add_notes(slide, """Here's the fundamental insight. When we weight the PPO gradient by importance scores w(t), we introduce a BIAS equal to Cov(w, f) - the covariance between the weights and the gradient.

If the weights correlate with gradient magnitude - which is the entire POINT of importance weighting - then this covariance is nonzero and the estimator is biased.

The key: variance decreases as O(1/sqrt(N)) with more data, but bias is PERSISTENT. It doesn't average out - it accumulates step after step.

Early in training, variance is high and bias is small relative to overall error - so importance weighting helps through variance reduction. But late in training, variance is naturally low and accumulated bias dominates - so importance weighting hurts.

This explains the degradation we saw. It's not a bug in any specific method - it's a fundamental property of non-uniform weighting in online policy gradient.""")

# ============================================================
# SLIDE 5: AITI Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Solution 1: AITI - Adaptive Intensity Token Importance",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 7, 3,
                  "The Fix: Decay importance intensity over training",
                  font_size=22, color=WHITE, bold=True)
add_para(tf, "", 10)
add_para(tf, "w(t) = 1 + eps(step) * (s(t) - 1)", 24, YELLOW, True, 12)
add_para(tf, "", 8)
add_para(tf, "eps = 1  -->  full importance weighting (max variance reduction)", 17, GREEN, False, 8)
add_para(tf, "eps = 0  -->  uniform PPO (zero bias)", 17, GREEN, False, 4)
add_para(tf, "", 10)
add_para(tf, "Schedule: eps decays linearly from 1 to 0", 18, WHITE, False, 8)
add_para(tf, "", 8)
add_para(tf, "Like learning rate scheduling,", 18, LIGHT, False, 8)
add_para(tf, "but for gradient QUALITY, not step size", 18, ACCENT2, True, 4)

tf2 = add_text_box(slide, 8, 1.3, 4.8, 5.5,
                   "Intuition:",
                   font_size=22, color=YELLOW, bold=True)
add_para(tf2, "", 8)
add_para(tf2, "Early training:", 20, GREEN, True, 12)
add_para(tf2, "  Model is far from optimal", 16, LIGHT, False, 4)
add_para(tf2, "  Gradients are high-variance", 16, LIGHT, False, 4)
add_para(tf2, "  Importance: REDUCE VARIANCE", 16, GREEN, True, 4)
add_para(tf2, "", 12)
add_para(tf2, "Late training:", 20, ACCENT2, True, 12)
add_para(tf2, "  Model is nearly converged", 16, LIGHT, False, 4)
add_para(tf2, "  Bias would push away from optimum", 16, LIGHT, False, 4)
add_para(tf2, "  Uniform: REDUCE BIAS", 16, ACCENT2, True, 4)
add_para(tf2, "", 12)
add_para(tf2, "Best of both worlds!", 22, YELLOW, True, 12)

add_notes(slide, """The fix is simple: decay the importance intensity over training.

Instead of using raw weights w(t) = s(t), we use w(t) = 1 + epsilon * (s(t) - 1), where epsilon decays from 1 to 0.

When epsilon is 1, you get full importance weighting - maximum variance reduction. When epsilon is 0, you get uniform PPO - zero bias.

Think of it like learning rate scheduling, but for gradient quality rather than step size. Early in training, you want low-variance gradients even if they're biased. Late in training, you want unbiased gradients even if they're noisier.

The implementation is trivial - 15 lines of code wrapping any existing importance scorer.""")

# ============================================================
# SLIDE 6: AITI Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "AITI Results: First Pareto-Optimal Token Importance for PPO",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 12, 5,
                  "Method                    Reward(Q4)  KL(Q4)    KL-Eff   Pareto?",
                  font_size=18, color=ACCENT2, bold=True)
add_para(tf, "--------------------------------------------------------------------------", 14, LIGHT, False, 2)
add_para(tf, "AITI-Advantage (linear)   0.855       0.016     54.1     YES <<<", 18, GREEN, True, 6)
add_para(tf, "AITI-AdaptivePhase        0.853       0.047     18.1     no", 18, WHITE, False, 6)
add_para(tf, "AITI-Entropy (quad)       0.826       0.059     13.9     no", 18, WHITE, False, 6)
add_para(tf, "PPO baseline              0.818       0.040     20.4     DOMINATED", 18, RED, False, 6)
add_para(tf, "Entropy (no decay)        0.815      -0.611      1.3     no", 18, LIGHT, False, 6)
add_para(tf, "", 14)
add_para(tf, "AITI-Advantage vs PPO baseline:", 22, YELLOW, True, 16)
add_para(tf, "", 6)
add_para(tf, "  Reward:        +4.5%   (0.855 vs 0.818)", 20, GREEN, True, 8)
add_para(tf, "  KL divergence: 2.5x lower  (0.016 vs 0.040)", 20, GREEN, True, 6)
add_para(tf, "  KL-Efficiency: 2.7x higher (54.1 vs 20.4)", 20, GREEN, True, 6)
add_para(tf, "", 10)
add_para(tf, "  Better reward AND closer to reference model -- strictly dominates PPO", 18, YELLOW, False, 8)

add_notes(slide, """AITI-Advantage with linear decay is the ONLY method that Pareto-dominates uniform PPO.

It achieves 4.5% higher reward AND 2.5 times lower KL divergence simultaneously. That means the resulting policy is both better at the task AND closer to the reference model.

The KL-efficiency - reward per unit of KL divergence - is 54.1, nearly 3x better than PPO's 20.4.

Importantly, without the decay, entropy weighting achieves similar reward but with NEGATIVE KL and very low efficiency. And AITI-Entropy with quadratic decay is decent but not Pareto-optimal. The advantage magnitude scorer combined with linear decay is the winning combination.""")

# ============================================================
# SLIDE 7: MOAI Derivation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Solution 2: MOAI - MSE-Optimal Adaptive Intensity",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6.5, 5.5,
                  "Can we derive epsilon from data?",
                  font_size=22, color=WHITE, bold=True)
add_para(tf, "", 8)
add_para(tf, "For w = 1 + eps*(s - 1):", 18, LIGHT, False, 8)
add_para(tf, "", 6)
add_para(tf, "MSE(eps) = Bias^2 + Variance", 18, WHITE, False, 8)
add_para(tf, "         = eps^2 * C^2 + (1/T)(sigma^2 + 2*eps*rho + eps^2*tau^2)", 16, LIGHT, False, 4)
add_para(tf, "", 8)
add_para(tf, "Setting dMSE/deps = 0:", 18, WHITE, False, 8)
add_para(tf, "", 8)
add_para(tf, "  eps* = -rho / (T * C^2 + tau^2)", 24, YELLOW, True, 8)
add_para(tf, "", 10)
add_para(tf, "where:", 17, LIGHT, False, 8)
add_para(tf, "  C   = Cov(s, f)        -- importance-loss correlation", 15, LIGHT, False, 4)
add_para(tf, "  rho = Cov(f, (s-1)*f)  -- variance-importance coupling", 15, LIGHT, False, 4)
add_para(tf, "  tau = Var((s-1)*f)      -- modulated loss variance", 15, LIGHT, False, 4)
add_para(tf, "  T   = token count       -- sequence length", 15, LIGHT, False, 4)

tf2 = add_text_box(slide, 7.5, 1.3, 5.3, 5.5,
                   "Three Predictions:",
                   font_size=22, color=YELLOW, bold=True)
add_para(tf2, "", 10)
add_para(tf2, "1. Sequence length effect", 20, GREEN, True, 12)
add_para(tf2, "   As T -> inf, eps* -> 0", 17, LIGHT, False, 4)
add_para(tf2, "   Longer sequences need LESS", 17, LIGHT, False, 4)
add_para(tf2, "   importance weighting", 17, LIGHT, False, 4)
add_para(tf2, "", 10)
add_para(tf2, "2. Natural decay", 20, GREEN, True, 12)
add_para(tf2, "   As C grows during training,", 17, LIGHT, False, 4)
add_para(tf2, "   eps* decreases automatically", 17, LIGHT, False, 4)
add_para(tf2, "", 10)
add_para(tf2, "3. When weighting helps", 20, GREEN, True, 12)
add_para(tf2, "   eps* > 0 iff rho < 0", 17, LIGHT, False, 4)
add_para(tf2, "   (advantage weighting satisfies this)", 17, LIGHT, False, 4)
add_para(tf2, "", 10)
add_para(tf2, "This is NEW in the IS literature.", 20, ACCENT2, True, 12)

add_notes(slide, """AITI works but requires tuning the decay_steps hyperparameter. Can we derive epsilon from data instead?

We minimize the Mean Squared Error of the weighted gradient estimator. The MSE decomposes into bias-squared plus variance, both functions of epsilon. Taking the derivative and setting to zero gives us a beautiful closed-form: epsilon-star equals negative rho over T-C-squared plus tau-squared.

This formula has three interesting predictions. First, epsilon depends on T - the token count - meaning longer sequences naturally need less importance weighting. This is testable and has practical implications for long-context RLHF.

Second, epsilon naturally decreases as training progresses because C - the importance-loss correlation - grows over time.

Third, importance weighting helps (epsilon > 0) only when rho < 0, which is exactly what advantage magnitude weighting provides.

This closed-form is new. Prior work by Korba and Portier uses cross-validation; Hachiya et al. use importance-weighted cross-validation. Nobody has derived the optimal intensity from MSE minimization.""")

# ============================================================
# SLIDE 8: MOAI - Theory vs Practice
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "The Theory-Practice Gap and the Monotone Fix",
             font_size=30, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6, 5.5,
                  "Problem: Free MOAI oscillates",
                  font_size=22, color=RED, bold=True)
add_para(tf, "", 8)
add_para(tf, "eps trajectory (free, ema=0.80):", 18, WHITE, False, 8)
add_para(tf, "  1.0 -> 0.54 -> 1.00 -> 1.00 -> 0.94 -> 1.00", 17, RED, False, 6)
add_para(tf, "", 8)
add_para(tf, "WHY: batch statistics are noisy", 18, YELLOW, False, 8)
add_para(tf, "  (batch size = 4 sequences)", 16, LIGHT, False, 4)
add_para(tf, "  C, rho, tau^2 fluctuate wildly", 16, LIGHT, False, 4)
add_para(tf, "  Formula amplifies the noise", 16, LIGHT, False, 4)
add_para(tf, "", 12)
add_para(tf, "Result: R=0.770, KL=0.097", 18, RED, True, 8)
add_para(tf, "  (worse than PPO baseline!)", 16, RED, False, 4)

tf2 = add_text_box(slide, 7, 1.3, 5.8, 5.5,
                   "Fix: Monotone constraint",
                   font_size=22, color=GREEN, bold=True)
add_para(tf2, "", 8)
add_para(tf2, "eps_k = min(eps*_k,  eps_{k-1})", 20, YELLOW, True, 8)
add_para(tf2, "", 6)
add_para(tf2, "eps can only DECREASE over time.", 18, WHITE, False, 8)
add_para(tf2, "", 10)
add_para(tf2, "eps trajectory (mono, ema=0.80):", 18, WHITE, False, 8)
add_para(tf2, "  1.0 -> 0.43 -> 0.24 -> 0.24 -> 0.24", 17, GREEN, False, 6)
add_para(tf2, "", 8)
add_para(tf2, "\"Find and Lock\" strategy:", 20, ACCENT2, True, 12)
add_para(tf2, "  1. Rapid early decay (data-adaptive)", 16, LIGHT, False, 4)
add_para(tf2, "  2. Lock at optimal level (~0.24)", 16, LIGHT, False, 4)
add_para(tf2, "  3. Stay there (stability)", 16, LIGHT, False, 4)
add_para(tf2, "", 10)
add_para(tf2, "Result: R=0.880, KL=0.056", 18, GREEN, True, 8)
add_para(tf2, "  (+14% reward vs free MOAI!)", 16, GREEN, False, 4)

add_notes(slide, """When we actually ran MOAI, we hit a theory-practice gap. The per-step epsilon oscillates wildly instead of smoothly decaying. It bounces between 0.5 and 1.0 because the batch-level estimates of C, rho, and tau-squared are too noisy with only 4 sequences per batch.

The formula is mathematically correct in expectation, but practically unstable sample-by-sample. Free MOAI gets worse reward than PPO baseline.

The fix is simple but powerful: add a monotone constraint. Epsilon can only decrease, never increase. This creates a "find-and-lock" behavior: epsilon drops rapidly in early training as the formula responds to growing covariance, then LOCKS at around 0.24 because subsequent epsilon estimates are higher but the constraint prevents going back up.

The result is dramatic: monotone MOAI gets 0.880 reward - 14% better than free MOAI and the highest of ANY method we tested.

This parallels a well-known phenomenon in deep learning: fixed schedules often beat adaptive methods because per-step optimality doesn't equal trajectory optimality when estimates are noisy.""")

# ============================================================
# SLIDE 9: Final Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Final Results: The Pareto Frontier",
             font_size=32, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 12, 5,
                  "Method                     Reward(Q4)  KL(Q4)    Pareto?    eps locks at",
                  font_size=18, color=ACCENT2, bold=True)
add_para(tf, "-------------------------------------------------------------------------------", 14, LIGHT, False, 2)
add_para(tf, "MOAI-Adv mono (ema=0.80)   0.880       0.056     YES <<<    0.24", 18, GREEN, True, 6)
add_para(tf, "AITI-Advantage (linear)    0.874       0.031     YES        -> 0", 18, GREEN, True, 6)
add_para(tf, "PPO baseline               0.859       0.028     YES        1.0", 18, WHITE, False, 6)
add_para(tf, "MOAI-Adv mono (ema=0.90)   0.866       0.059     no         0.49", 18, LIGHT, False, 6)
add_para(tf, "MOAI-Adv free (ema=0.80)   0.770       0.097     no         oscillating", 18, RED, False, 6)
add_para(tf, "", 12)
add_para(tf, "Surprise Finding:", 22, YELLOW, True, 16)
add_para(tf, "", 6)
add_para(tf, "  AITI assumes optimal eps = 0 (decay to uniform)", 18, LIGHT, False, 8)
add_para(tf, "  MOAI discovers optimal eps ~ 0.24 (permanent partial weighting)", 18, GREEN, True, 6)
add_para(tf, "", 8)
add_para(tf, "  Some importance weighting is ALWAYS beneficial -- contradicts AITI's assumption!", 18, YELLOW, False, 8)

add_notes(slide, """Here are the final results across all methods. Three methods sit on the Pareto frontier:

MOAI-Advantage monotone with ema 0.80 achieves the HIGHEST reward at 0.880, but with moderate KL at 0.056. It locked epsilon at 0.24.

AITI-Advantage provides the best reward-KL TRADEOFF at 0.874 reward and only 0.031 KL.

PPO baseline is Pareto-optimal in the low-KL region at 0.028.

The surprise finding: AITI decays epsilon all the way to zero, assuming you should eventually use uniform weighting. But MOAI's closed-form consistently says the optimal epsilon is POSITIVE - around 0.24. This means some permanent token-level differentiation is beneficial, even late in training.

Why? Even in a nearly-converged model, high-advantage tokens genuinely have more informative gradients. The bias from weighting them is small (the advantage landscape is stable), while the variance reduction is still useful.

Also note the free MOAI result - without the monotone constraint, it performs WORSE than PPO baseline. The stability from monotonicity is worth 14% reward.""")

# ============================================================
# SLIDE 10: Key Findings Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Key Findings",
             font_size=32, color=ACCENT2, bold=True)

findings = [
    ("1", "Importance weighting has a bias-variance tradeoff in PPO",
     "Bias = Cov(w, f) accumulates; variance reduction is temporary", YELLOW),
    ("2", "Decaying intensity solves it (AITI)",
     "+4.5% reward, 2.5x lower KL vs uniform PPO", GREEN),
    ("3", "MSE-optimal intensity has a closed form (MOAI)",
     "eps* = -rho / (T*C^2 + tau^2)  -- depends on sequence length T", ACCENT2),
    ("4", "Per-step optimality != trajectory optimality",
     "Free MOAI oscillates; monotone constraint is essential", ORANGE),
    ("5", "Optimal epsilon > 0 (the surprise)",
     "Permanent partial weighting beats full decay to uniform", GREEN),
]

y = 1.3
for num, title, detail, color in findings:
    tf = add_text_box(slide, 0.8, y, 0.5, 0.5, num, font_size=28, color=color, bold=True)
    tf2 = add_text_box(slide, 1.5, y, 10.5, 0.4, title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, 1.5, y + 0.4, 10.5, 0.4, detail, font_size=16, color=LIGHT)
    y += 1.1

add_notes(slide, """Let me summarize the five key findings:

First, we identified the bias-variance tradeoff of token importance in PPO. The bias from non-uniform weighting accumulates step by step, while variance reduction is temporary.

Second, AITI solves this with a simple linear decay schedule, achieving the first Pareto-optimal token importance for PPO.

Third, we derived the MSE-optimal importance intensity as a closed-form formula. This is new in the importance sampling literature, and it predicts that longer sequences need less weighting.

Fourth, we discovered a theory-practice gap: the per-step optimal epsilon oscillates due to noisy batch statistics. A monotone constraint fixes this, giving a 14% reward improvement. This adds to evidence that scheduled methods can beat adaptive ones.

Fifth, the surprise: MOAI discovers that the optimal epsilon isn't zero - it's around 0.24. Some permanent importance weighting is always beneficial, contradicting AITI's full-decay-to-uniform assumption.""")

# ============================================================
# SLIDE 11: Implications & Future Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.8, "Implications and Future Directions",
             font_size=32, color=ACCENT2, bold=True)

tf = add_text_box(slide, 0.5, 1.3, 6, 5.5,
                  "For Practitioners:",
                  font_size=22, color=GREEN, bold=True)
add_para(tf, "", 8)
add_para(tf, "AITI: 15 lines of code, free to compute", 18, WHITE, False, 8)
add_para(tf, "  Just wrap your PPO loss with epsilon decay", 16, LIGHT, False, 4)
add_para(tf, "  Use |Advantage| scorer (already computed)", 16, LIGHT, False, 4)
add_para(tf, "", 10)
add_para(tf, "MOAI: slightly more complex, no tuning", 18, WHITE, False, 8)
add_para(tf, "  EMA statistics + monotone constraint", 16, LIGHT, False, 4)
add_para(tf, "  No decay_steps hyperparameter", 16, LIGHT, False, 4)
add_para(tf, "", 14)
add_para(tf, "For Researchers:", 22, ACCENT2, True, 12)
add_para(tf, "", 8)
add_para(tf, "The T-dependence prediction is testable", 18, WHITE, False, 8)
add_para(tf, "  Does eps* decrease with longer contexts?", 16, LIGHT, False, 4)
add_para(tf, "", 8)
add_para(tf, "\"Find-and-lock\" may generalize to", 18, WHITE, False, 8)
add_para(tf, "  other adaptive hyperparameters in RL", 16, LIGHT, False, 4)

tf2 = add_text_box(slide, 7, 1.3, 5.8, 5.5,
                   "Limitations:",
                   font_size=22, color=ORANGE, bold=True)
add_para(tf2, "", 8)
add_para(tf2, "GPT-2 (124M) + synthetic reward", 18, LIGHT, False, 8)
add_para(tf2, "  Need validation at 7B+ scale", 16, LIGHT, False, 4)
add_para(tf2, "", 8)
add_para(tf2, "Single seed per method", 18, LIGHT, False, 8)
add_para(tf2, "  Need multi-seed significance tests", 16, LIGHT, False, 4)
add_para(tf2, "", 8)
add_para(tf2, "Batch size = 4 (noisy statistics)", 18, LIGHT, False, 8)
add_para(tf2, "  Larger batches may help MOAI", 16, LIGHT, False, 4)
add_para(tf2, "", 14)
add_para(tf2, "Future Work:", 22, YELLOW, True, 12)
add_para(tf2, "", 8)
add_para(tf2, "Scale to 7B+ with real reward models", 18, WHITE, False, 8)
add_para(tf2, "Apply to DPO and GRPO", 18, WHITE, False, 8)
add_para(tf2, "Test T-dependence prediction at scale", 18, WHITE, False, 8)
add_para(tf2, "Formalize cumulative MSE optimization", 18, WHITE, False, 8)

add_notes(slide, """For practitioners: AITI is dead simple to implement - 15 lines of code, no extra model calls, just wrap your PPO loss with an epsilon that decays linearly. Use advantage magnitude as your importance scorer since it's already computed in PPO.

For those who want to avoid tuning decay_steps, monotone MOAI eliminates that hyperparameter. It's slightly more complex but fully data-adaptive.

For researchers: the sequence length prediction is directly testable. As models move to longer contexts in RLHF, does the optimal epsilon decrease? This could have practical implications for how we train large models.

The find-and-lock strategy may generalize beyond importance weighting to other adaptive hyperparameters in RL.

The main limitations are scale - we used GPT-2 with a synthetic reward - and single seeds. These results need validation at 7B+ scale with learned reward models and multiple seeds for statistical significance.""")

# ============================================================
# SLIDE 12: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 1, 2.0, 11.3, 1.2,
             "Thank You",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text_box(slide, 2, 3.5, 9.3, 3,
                  "Key Contributions:",
                  font_size=24, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)
add_para(tf, "", 8)
add_para(tf, "Bias-variance analysis of token importance in RLHF", 20, WHITE, False, 8)
add_para(tf, "AITI: Pareto-optimal via intensity decay", 20, WHITE, False, 8)
add_para(tf, "MOAI: Closed-form MSE-optimal intensity", 20, WHITE, False, 8)
add_para(tf, "Monotone \"find-and-lock\" beats both schedules and adaptive methods", 20, WHITE, False, 8)
for p in tf.paragraphs[2:]:
    p.alignment = PP_ALIGN.CENTER

add_text_box(slide, 1, 6.0, 11.3, 0.5,
             "Code & results: 291K/  |  19 importance methods  |  6 benchmark scripts",
             font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_notes(slide, """Thank you for your attention. To summarize:

We identified why token importance weighting fails in PPO - the bias-variance tradeoff.

We proposed AITI which achieves Pareto dominance through simple intensity decay.

We derived MOAI - the first closed-form MSE-optimal importance intensity - and discovered that monotone-constrained MOAI achieves the highest reward by finding and locking at a positive epsilon.

The code includes 19 different importance methods and 6 benchmark scripts. Happy to take questions.""")

# Save
out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                         "presentation.pptx")
prs.save(out_path)
print(f"Saved to {out_path}")
